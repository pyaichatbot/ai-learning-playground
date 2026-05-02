from pathlib import Path
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
COMIC_DIR = ROOT / "comic"
SRC_DIR = ROOT / "comic_enterprise/assets_ai"
WORK = COMIC_DIR / "enterprise_v2_work"
PREP = WORK / "prepared"
FINAL = COMIC_DIR / "ai-migration-comic-enterprise-v2.pptx"

INK = (23, 31, 42)
PAPER = (244, 241, 234)
WARM = (190, 116, 72)
TEAL = (44, 124, 130)
BLUE = (54, 91, 150)
BURGUNDY = (124, 57, 76)
GOLD = (205, 154, 78)
MUTED = (90, 99, 112)
WHITE = (255, 255, 255)


def font_size(text, base=9.5):
    if len(text) > 110:
        return base - 1.3
    if len(text) > 80:
        return base - 0.8
    return base


def prepare_images():
    PREP.mkdir(parents=True, exist_ok=True)
    paths = sorted(SRC_DIR.glob("*.png"))
    out = []
    for i, p in enumerate(paths, 1):
        im = Image.open(p).convert("RGB")
        w, h = im.size
        # Convert generated 3:2 panels into 16:9 without distortion. Keep top area
        # because the AI often reserves clean space there for comic dialogue.
        target_h = int(w * 9 / 16)
        y0 = max(0, min(72, h - target_h))
        im = im.crop((0, y0, w, y0 + target_h)).resize((1600, 900), Image.Resampling.LANCZOS)
        # Subtle sharpening and film grain for a consistent printed-comic finish.
        im = im.filter(ImageFilter.UnsharpMask(radius=1.1, percent=125, threshold=4))
        overlay = Image.new("RGBA", im.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(overlay)
        for y in range(0, 900, 22):
            for x in range((y // 22 % 2) * 11, 1600, 22):
                d.ellipse((x, y, x + 3, y + 3), fill=(0, 0, 0, 12))
        im = Image.alpha_composite(im.convert("RGBA"), overlay).convert("RGB")
        op = PREP / f"panel_{i:02d}.jpg"
        im.save(op, quality=94)
        out.append(op)
    return out


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    return box


def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*color)


def title_bar(slide, kicker, title):
    add_text(slide, 0.45, 0.22, 2.4, 0.28, kicker.upper(), 8.5, True, WARM)
    add_text(slide, 0.45, 0.48, 9.8, 0.42, title, 22, True, INK)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.95), Inches(12.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*INK)
    line.line.fill.background()


def add_panel(slide, img, x, y, w, h, caption, bubbles, artifact=None):
    slide.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(w), Inches(h))
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    border.fill.background()
    border.line.color.rgb = RGBColor(*INK)
    border.line.width = Pt(2.5)
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + h - 0.42), Inches(w), Inches(0.42))
    cap.fill.solid()
    cap.fill.fore_color.rgb = RGBColor(*INK)
    cap.line.fill.background()
    add_text(slide, x + 0.1, y + h - 0.38, w - 0.2, 0.28, caption, 8.2, True, WHITE, PP_ALIGN.CENTER)
    if artifact:
        add_artifact_chip(slide, x + 0.12, y + 0.1, artifact, BLUE)
    for bx, by, bw, bh, text, accent in bubbles:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, Inches(x + bx * w), Inches(y + by * h), Inches(bw * w), Inches(bh * h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(255, 255, 255)
        sh.line.color.rgb = RGBColor(*accent)
        sh.line.width = Pt(1.6)
        tf = sh.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Aptos"
        r.font.bold = True
        r.font.size = Pt(font_size(text))
        r.font.color.rgb = RGBColor(*INK)


def add_pill(slide, x, y, w, text, fill):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    p.fill.solid()
    p.fill.fore_color.rgb = RGBColor(*fill)
    p.line.fill.background()
    add_text(slide, x + 0.06, y + 0.05, w - 0.12, 0.24, text, 8.5, True, WHITE, PP_ALIGN.CENTER)


def add_artifact_chip(slide, x, y, text, fill=BLUE):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(0.28))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(*fill)
    chip.line.fill.background()
    add_text(slide, x + 0.08, y + 0.035, 2.08, 0.18, text, 7.0, True, WHITE, PP_ALIGN.CENTER)


def add_yaml_card(slide, x, y, w, h, title, yaml_text):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(*INK)
    card.line.width = Pt(1.6)
    add_text(slide, x + 0.18, y + 0.16, w - 0.36, 0.25, title, 11, True, BLUE)
    add_text(slide, x + 0.22, y + 0.58, w - 0.44, h - 0.76, yaml_text, 8.3, False, INK)


def build(paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    s = prs.slides.add_slide(blank)
    bg(s)
    s.shapes.add_picture(str(paths[0]), Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    veil = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    veil.fill.solid()
    veil.fill.fore_color.rgb = RGBColor(8, 16, 26)
    veil.fill.transparency = 18
    veil.line.fill.background()
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(0.68), Inches(5.9), Inches(5.78))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(246, 243, 235)
    panel.fill.transparency = 4
    panel.line.color.rgb = RGBColor(*INK)
    panel.line.width = Pt(2)
    add_text(s, 0.92, 1.0, 4.8, 0.28, "AI-SDLC FRAMEWORK", 9.5, True, WARM)
    add_text(s, 0.9, 1.35, 4.95, 1.35, "Migration Walkthrough\nComic Tutorial", 33, True, INK)
    add_text(s, 0.92, 3.05, 4.8, 0.75, "An enterprise architect walkthrough: how teams move from legacy pressure to governed AI-assisted delivery.", 14, False, INK)
    add_pill(s, 0.92, 4.25, 1.3, "Govern", BURGUNDY)
    add_pill(s, 2.35, 4.25, 1.25, "Architect", TEAL)
    add_pill(s, 3.75, 4.25, 1.16, "Agents", BLUE)
    add_pill(s, 5.05, 4.25, 0.92, "Scale", GOLD)
    add_text(s, 0.92, 5.62, 4.95, 0.44, "Presenter stance: framework first, agents second, evidence always.", 11, True, INK)

    # Cast and premise
    s = prs.slides.add_slide(blank)
    bg(s)
    title_bar(s, "Architect setup", "The story starts with tension, not a tool demo")
    add_panel(s, paths[0], 0.55, 1.22, 7.15, 5.35, "Decision: enter the AI SDLC path before selecting any coding agent.", [
        (0.05, 0.08, 0.31, 0.16, "Maya: This is now a delivery risk.", BLUE),
        (0.38, 0.08, 0.30, 0.16, "Omar: Then audit cannot be optional.", GOLD),
        (0.67, 0.14, 0.28, 0.16, "Rina: No agent before governance.", BURGUNDY),
    ], "G0 -> project-manifest.yml")
    add_text(s, 8.12, 1.28, 4.45, 0.38, "How to present it", 20, True, INK)
    cast = [("Maya", "Product owner", BLUE), ("Rina", "AI architecture lead", BURGUNDY), ("Dev", "Engineering lead", TEAL), ("Omar", "Risk stakeholder", GOLD)]
    for i, (name, role, color) in enumerate(cast):
        y = 1.92 + i * 0.72
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.15), Inches(y), Inches(4.45), Inches(0.52))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*INK)
        add_text(s, 8.35, y + 0.08, 0.95, 0.24, name, 11, True, color)
        add_text(s, 9.35, y + 0.08, 2.8, 0.24, role, 10.5, False, INK)
    add_text(s, 8.15, 5.0, 4.35, 0.9, "As the enterprise architect, narrate the operating model: the team can disagree, but the artifact settles the next gate.", 15, True, INK)
    add_text(s, 8.15, 6.15, 4.35, 0.45, "Core line: The framework is not an abstraction layer. The manifest is the control plane.", 12.5, True, BURGUNDY)

    # Act 1
    s = prs.slides.add_slide(blank)
    bg(s)
    title_bar(s, "Act 1", "Controlled intake before agent adoption")
    pos = [(0.45, 1.18), (6.9, 1.18), (0.45, 4.1), (6.9, 4.1)]
    specs = [
        (paths[0], "Decision: treat migration delay as an enterprise risk.", "G0 -> intake", [("Maya: Six months per feature is the problem.", BLUE), ("Dev: Every change still touches the monolith.", TEAL)]),
        (paths[1], "Decision: block agent usage until controls exist.", "G2 -> risk", [("Dev: Can Copilot clear the backlog?", TEAL), ("Omar: Not without audit trails.", GOLD)]),
        (paths[2], "Decision: let the manifest determine the next gate.", "G0/G2 -> manifest", [("Rina: We classify risk first.", BURGUNDY), ("Maya: So the artifact decides next.", BLUE)]),
        (paths[3], "Decision: migrate by domain slice, not big bang.", "G3 -> architecture", [("Rina: Map capabilities and data flows.", BURGUNDY), ("Dev: Then pilot one thin slice.", TEAL)]),
    ]
    for (x, y), (img, cap, artifact, pairs) in zip(pos, specs):
        add_panel(s, img, x, y, 5.98, 2.64, cap, [(0.04, 0.07, 0.41, 0.18, pairs[0][0], pairs[0][1]), (0.54, 0.07, 0.40, 0.18, pairs[1][0], pairs[1][1])], artifact)

    # Act 2
    s = prs.slides.add_slide(blank)
    bg(s)
    title_bar(s, "Act 2", "Guardrails, subagents, and evidence gates")
    specs = [
        (paths[4], "Decision: guardrails travel with every agent task.", "G3 -> ADR links", [("Omar: Who owns traceability?", GOLD), ("Rina: The manifest links every decision.", BURGUNDY)]),
        (paths[5], "Decision: agents vary, SDLC contract stays fixed.", "G4 -> agent model", [("Dev: Codex here, Claude there?", TEAL), ("Maya: Same gates, different assistants.", BLUE)]),
        (paths[6], "Decision: no merge without evidence.", "G5 -> evidence", [("Rina: Tests and risk notes first.", BURGUNDY), ("Omar: The team approves, not the tool.", GOLD)]),
        (paths[7], "Decision: scale only what becomes measurable.", "G7 -> scale", [("Maya: Now value is visible each sprint.", BLUE), ("Dev: And reuse is safe.", TEAL)]),
    ]
    for (x, y), (img, cap, artifact, pairs) in zip(pos, specs):
        add_panel(s, img, x, y, 5.98, 2.64, cap, [(0.04, 0.07, 0.41, 0.19, pairs[0][0], pairs[0][1]), (0.54, 0.07, 0.40, 0.19, pairs[1][0], pairs[1][1])], artifact)

    # Manifest gate slide
    s = prs.slides.add_slide(blank)
    bg(s)
    title_bar(s, "Manifest gate", "Make the framework tangible at every decision")
    add_panel(s, paths[2], 0.55, 1.25, 6.35, 4.95, "Decision: the team does not move gates until the manifest is updated.", [
        (0.06, 0.08, 0.38, 0.18, "Rina: This is not abstraction.", BURGUNDY),
        (0.53, 0.08, 0.39, 0.18, "Omar: Show me the control record.", GOLD),
    ], "control plane")
    manifest_yaml = """project:
  scenario: migration
  owner: product-platform
  criticality: high
ai_sdlc:
  onboarding_mode: onboard-lite
  current_gate: G2-risk
risk:
  classification: high
  human_approval_required: true
architecture:
  decision_records:
    - adr/001-domain-slice.md
agents:
  allowed: [codex, claude, copilot]
evidence:
  required: [tests, risk_notes, arch_review]"""
    add_yaml_card(s, 7.35, 1.25, 5.25, 4.95, "project-manifest.yml", manifest_yaml)
    add_text(s, 0.85, 6.55, 11.8, 0.34, "Presenter cue: point to the YAML, not the process name. The manifest is how governance becomes executable.", 13, True, INK, PP_ALIGN.CENTER)

    # Decision spine
    s = prs.slides.add_slide(blank)
    bg(s, (238, 240, 236))
    title_bar(s, "Architect walkthrough", "Use the comic as a guided decision conversation")
    add_text(s, 0.72, 1.28, 11.6, 0.36, "For each panel, ask: What decision is being made? Which artifact enforces it? What proof unlocks the next gate?", 15, True, INK, PP_ALIGN.CENTER)
    steps = [
        ("Problem", "Legacy delay and unclear ownership", BURGUNDY),
        ("Decision", "Classify risk and choose the migration slice", GOLD),
        ("Artifact", "Manifest, knowledge, ADRs, subagents", BLUE),
        ("Gate", "Evidence before merge and release", TEAL),
        ("Outcome", "Measured modernization", (82, 137, 94)),
    ]
    for i, (head, body, col) in enumerate(steps):
        x = 0.65 + i * 2.52
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.24), Inches(2.05), Inches(2.25))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*INK)
        card.line.width = Pt(2)
        add_text(s, x + 0.2, 2.55, 1.65, 0.28, head, 14, True, col, PP_ALIGN.CENTER)
        add_text(s, x + 0.2, 3.05, 1.65, 0.7, body, 10.5, False, INK, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.08), Inches(3.03), Inches(0.38), Inches(0.48))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(*INK)
            arr.line.fill.background()
    add_text(s, 1.15, 5.3, 11.0, 0.55, "The story lands when the audience sees one thing clearly: the framework turns tool usage into governed SDLC decisions.", 18, True, INK, PP_ALIGN.CENTER)
    add_pill(s, 3.52, 6.18, 1.4, "Claude", BURGUNDY)
    add_pill(s, 5.18, 6.18, 1.3, "Codex", TEAL)
    add_pill(s, 6.75, 6.18, 1.85, "GitHub Copilot", BLUE)
    add_pill(s, 8.86, 6.18, 1.42, "Future agents", GOLD)

    # POC slide
    s = prs.slides.add_slide(blank)
    bg(s)
    title_bar(s, "Operating model", "When teams should engage the AI SDLC framework POC")
    add_text(s, 0.75, 1.32, 6.05, 0.62, "Use this closing slide to position ownership operationally, not personally: you are the entry point for consistent framework enablement.", 16, True, INK)
    services = [
        ("Classify", "Use-case and risk class"),
        ("Onboard", "Manifest + knowledge setup"),
        ("Architect", "ADRs, boundaries, data flows"),
        ("Enable", "Prompts, commands, subagents"),
        ("Evidence", "Tests, reviews, release proof"),
    ]
    for i, (head, body) in enumerate(services):
        y = 2.25 + i * 0.72
        add_pill(s, 0.85, y, 1.35, head, [BURGUNDY, BLUE, TEAL, GOLD, MUTED][i])
        add_text(s, 2.38, y + 0.02, 4.25, 0.28, body, 12.5, True, INK)
    add_yaml_card(s, 7.25, 1.55, 5.15, 4.8, "POC handoff contract", """before teams start:
  engage: ai-sdlc-framework-poc
  bring:
    - problem statement
    - target system
    - data sensitivity
    - intended agents
    - release path
framework_owner_provides:
  - onboarding path
  - manifest readiness
  - governance alignment
  - artifact templates
  - evidence gate review""")
    add_text(s, 0.9, 6.7, 11.6, 0.34, "Suggested close: 'The comic is the onboarding layer. The asset is the operating model behind it.'", 13, True, BURGUNDY, PP_ALIGN.CENTER)

    prs.save(FINAL)


def main():
    paths = prepare_images()
    build(paths)
    print(FINAL)


if __name__ == "__main__":
    main()
